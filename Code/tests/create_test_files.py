"""
Script to create sample PPTX and PDF test files for testing.
Run this script to generate test data files.
"""

from pathlib import Path


def create_sample_pptx(output_path: Path):
    """Create a sample PowerPoint file for testing."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()

        # Slide 1: Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "Introduction to Data Structures"
        subtitle.text = "A comprehensive guide to fundamental data structures"

        # Slide 2: Overview
        bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
        shapes = bullet_slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "What are Data Structures?"
        tf = body_shape.text_frame
        tf.text = "Data structures are specialized formats for organizing and storing data"
        p = tf.add_paragraph()
        p.text = "They enable efficient access and modification of data"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "Different structures are suited for different use cases"
        p.level = 1

        # Slide 3: Arrays
        array_slide = prs.slides.add_slide(prs.slide_layouts[1])
        shapes = array_slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Arrays"
        tf = body_shape.text_frame
        tf.text = "Arrays store elements in contiguous memory locations"
        p = tf.add_paragraph()
        p.text = "Access time: O(1) - constant time access by index"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "Insertion/Deletion: O(n) - may require shifting elements"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "Best for: Fixed-size collections with frequent access"
        p.level = 1

        # Slide 4: Linked Lists
        list_slide = prs.slides.add_slide(prs.slide_layouts[1])
        shapes = list_slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Linked Lists"
        tf = body_shape.text_frame
        tf.text = "Linked lists store elements with pointers to next elements"
        p = tf.add_paragraph()
        p.text = "Access time: O(n) - must traverse from head"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "Insertion/Deletion: O(1) - when position is known"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "Best for: Dynamic size, frequent insertions/deletions"
        p.level = 1

        # Slide 5: Stacks and Queues
        sq_slide = prs.slides.add_slide(prs.slide_layouts[1])
        shapes = sq_slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Stacks and Queues"
        tf = body_shape.text_frame
        tf.text = "Stack: LIFO (Last In, First Out) - like a stack of plates"
        p = tf.add_paragraph()
        p.text = "Operations: push, pop, peek - all O(1)"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "Queue: FIFO (First In, First Out) - like a line"
        p = tf.add_paragraph()
        p.text = "Operations: enqueue, dequeue, front - all O(1)"
        p.level = 1

        # Slide 6: Trees
        tree_slide = prs.slides.add_slide(prs.slide_layouts[1])
        shapes = tree_slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Trees"
        tf = body_shape.text_frame
        tf.text = "Trees are hierarchical structures with nodes and edges"
        p = tf.add_paragraph()
        p.text = "Binary Search Tree: O(log n) search, insert, delete"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "Balanced trees (AVL, Red-Black) maintain O(log n) height"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "Use cases: File systems, databases, decision trees"
        p.level = 1

        # Save the presentation
        prs.save(str(output_path))
        return True
    except ImportError:
        print("python-pptx not installed, skipping PPTX generation")
        return False


def create_sample_pdf(output_path: Path):
    """Create a sample PDF file for testing."""
    try:
        # Use reportlab if available, otherwise create a minimal PDF manually
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch

            doc = SimpleDocTemplate(str(output_path), pagesize=letter)
            styles = getSampleStyleSheet()
            story = []

            # Title
            title_style = ParagraphStyle(
                "CustomTitle",
                parent=styles["Heading1"],
                fontSize=24,
                spaceAfter=30,
            )
            story.append(Paragraph("Introduction to Algorithms", title_style))
            story.append(Spacer(1, 12))

            # Introduction
            story.append(Paragraph("Overview", styles["Heading2"]))
            story.append(
                Paragraph(
                    "An algorithm is a step-by-step procedure for solving a problem or "
                    "accomplishing a task. Algorithms are fundamental to computer science "
                    "and programming.",
                    styles["Normal"],
                )
            )
            story.append(Spacer(1, 12))

            # Sorting Algorithms
            story.append(Paragraph("Sorting Algorithms", styles["Heading2"]))
            story.append(
                Paragraph(
                    "Sorting algorithms arrange elements in a specific order. Common "
                    "sorting algorithms include:",
                    styles["Normal"],
                )
            )
            story.append(Spacer(1, 6))
            story.append(
                Paragraph(
                    "• Bubble Sort: O(n²) - Simple but inefficient for large datasets",
                    styles["Normal"],
                )
            )
            story.append(
                Paragraph(
                    "• Quick Sort: O(n log n) average - Divide and conquer approach",
                    styles["Normal"],
                )
            )
            story.append(
                Paragraph(
                    "• Merge Sort: O(n log n) - Stable, predictable performance",
                    styles["Normal"],
                )
            )
            story.append(
                Paragraph(
                    "• Heap Sort: O(n log n) - Uses heap data structure",
                    styles["Normal"],
                )
            )
            story.append(Spacer(1, 12))

            # Searching Algorithms
            story.append(Paragraph("Searching Algorithms", styles["Heading2"]))
            story.append(
                Paragraph(
                    "Searching algorithms find specific elements within a data structure:",
                    styles["Normal"],
                )
            )
            story.append(Spacer(1, 6))
            story.append(
                Paragraph(
                    "• Linear Search: O(n) - Checks each element sequentially",
                    styles["Normal"],
                )
            )
            story.append(
                Paragraph(
                    "• Binary Search: O(log n) - Requires sorted data, divides in half",
                    styles["Normal"],
                )
            )
            story.append(Spacer(1, 12))

            # Graph Algorithms
            story.append(Paragraph("Graph Algorithms", styles["Heading2"]))
            story.append(
                Paragraph(
                    "Graph algorithms solve problems on networks and relationships:",
                    styles["Normal"],
                )
            )
            story.append(Spacer(1, 6))
            story.append(
                Paragraph(
                    "• BFS (Breadth-First Search): Explores neighbors first",
                    styles["Normal"],
                )
            )
            story.append(
                Paragraph(
                    "• DFS (Depth-First Search): Explores as far as possible first",
                    styles["Normal"],
                )
            )
            story.append(
                Paragraph(
                    "• Dijkstra's Algorithm: Finds shortest path in weighted graphs",
                    styles["Normal"],
                )
            )
            story.append(Spacer(1, 12))

            # Complexity Analysis
            story.append(Paragraph("Algorithm Complexity", styles["Heading2"]))
            story.append(
                Paragraph(
                    "Big O notation describes the upper bound of an algorithm's time or "
                    "space complexity. Common complexities from fastest to slowest:",
                    styles["Normal"],
                )
            )
            story.append(Spacer(1, 6))
            story.append(Paragraph("• O(1) - Constant time", styles["Normal"]))
            story.append(Paragraph("• O(log n) - Logarithmic time", styles["Normal"]))
            story.append(Paragraph("• O(n) - Linear time", styles["Normal"]))
            story.append(Paragraph("• O(n log n) - Linearithmic time", styles["Normal"]))
            story.append(Paragraph("• O(n²) - Quadratic time", styles["Normal"]))
            story.append(Paragraph("• O(2^n) - Exponential time", styles["Normal"]))

            doc.build(story)
            return True

        except ImportError:
            # Fallback: Create a minimal PDF manually
            # This creates a valid PDF with basic text
            content = b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj
4 0 obj
<< /Length 200 >>
stream
BT
/F1 24 Tf
50 700 Td
(Introduction to Algorithms) Tj
/F1 12 Tf
0 -40 Td
(This is a sample PDF document for testing.) Tj
0 -20 Td
(It contains information about algorithms and data structures.) Tj
0 -20 Td
(Topics include sorting, searching, and graph algorithms.) Tj
ET
endstream
endobj
5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj
xref
0 6
0000000000 65535 f 
0000000009 00000 n 
0000000058 00000 n 
0000000115 00000 n 
0000000266 00000 n 
0000000518 00000 n 
trailer
<< /Size 6 /Root 1 0 R >>
startxref
595
%%EOF
"""
            with open(output_path, "wb") as f:
                f.write(content)
            return True

    except Exception as e:
        print(f"Error creating PDF: {e}")
        return False


def main():
    """Generate all sample test files."""
    test_data_dir = Path(__file__).parent / "test_data"
    test_data_dir.mkdir(exist_ok=True)

    print("Creating sample test files...")

    # Create PPTX
    pptx_path = test_data_dir / "sample.pptx"
    if create_sample_pptx(pptx_path):
        print(f"✓ Created: {pptx_path}")
    else:
        print(f"✗ Failed to create PPTX")

    # Create PDF
    pdf_path = test_data_dir / "sample.pdf"
    if create_sample_pdf(pdf_path):
        print(f"✓ Created: {pdf_path}")
    else:
        print(f"✗ Failed to create PDF")

    print("\nDone!")


if __name__ == "__main__":
    main()
