"""
PowerPoint (.pptx) extractor.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from study_guide.ingestion.extractors.base import BaseExtractor, ExtractionResult


class PPTXExtractor(BaseExtractor):
    """Extract text from PowerPoint files."""

    def supported_extensions(self) -> set[str]:
        return {".pptx"}

    def extract(self, filepath: Path | str) -> ExtractionResult:
        """
        Extract text from a PowerPoint file.

        Extracts text from:
        - Slide titles
        - Text boxes and shapes
        - Tables
        - Notes
        """
        try:
            filepath = self._validate_file(filepath)
            prs = Presentation(str(filepath))

            all_text: list[str] = []
            title = None
            slide_count = 0

            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_count += 1
                slide_text: list[str] = []

                # Add slide marker
                slide_text.append(f"\n--- Slide {slide_idx} ---\n")

                for shape in slide.shapes:
                    # Extract text from shapes with text frames
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                # Check if this is the title of the first slide
                                if slide_idx == 1 and title is None and shape.is_placeholder:
                                    placeholder_type = shape.placeholder_format.type
                                    # Title placeholder types
                                    if placeholder_type in (1, 3):  # TITLE, CENTER_TITLE
                                        title = text
                                slide_text.append(text)

                    # Extract text from tables
                    if shape.has_table:
                        table_text = self._extract_table(shape.table)
                        if table_text:
                            slide_text.append(table_text)

                # Extract notes if present
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text.strip():
                        slide_text.append(f"\nNotes: {notes_frame.text.strip()}")

                all_text.extend(slide_text)

            combined_text = "\n".join(all_text)

            return ExtractionResult(
                text=combined_text,
                title=title or filepath.stem,
                metadata={
                    "slide_count": slide_count,
                    "file_type": "pptx",
                },
                success=True,
            )

        except Exception as e:
            return ExtractionResult(
                text="",
                success=False,
                error=str(e),
            )

    def _extract_table(self, table) -> str:
        """Extract text from a table."""
        rows: list[str] = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    cells.append(cell_text)
            if cells:
                rows.append(" | ".join(cells))
        return "\n".join(rows)
